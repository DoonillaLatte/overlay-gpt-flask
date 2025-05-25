import base64
from io import BytesIO
from pptx import Presentation
import logging
from typing import Dict, Any, List, Optional

logger = logging.getLogger(__name__)

class PowerPointService:
    @staticmethod
    def read_pptx_data(base64_data: str) -> Optional[Presentation]:
        """
        base64로 인코딩된 PowerPoint 데이터를 읽어 Presentation 객체로 변환합니다.
        
        Args:
            base64_data (str): base64로 인코딩된 PowerPoint 파일 데이터
            
        Returns:
            Presentation: PowerPoint 프레젠테이션 객체
            None: 에러 발생 시
        """
        try:
            # base64 데이터를 디코딩
            pptx_data = base64.b64decode(base64_data)
            
            # BytesIO 객체로 변환
            pptx_file = BytesIO(pptx_data)
            
            # PowerPoint 파일 읽기
            presentation = Presentation(pptx_file)
            
            return presentation
        except Exception as e:
            logger.error(f"PowerPoint 파일 읽기 오류: {str(e)}")
            return None
    
    @staticmethod
    def get_pptx_summary(presentation: Presentation) -> Optional[Dict[str, Any]]:
        """
        PowerPoint 프레젠테이션의 기본적인 요약 정보를 반환합니다.
        
        Args:
            presentation (Presentation): PowerPoint 프레젠테이션 객체
            
        Returns:
            dict: 데이터 요약 정보
        """
        if presentation is None:
            return None
            
        try:
            # 슬라이드 스타일 정보 수집
            slides_info = []
            for slide in presentation.slides:
                slide_info = {
                    'layout_name': slide.slide_layout.name,
                    'shapes_count': len(slide.shapes),
                    'shapes': []
                }
                
                # 각 도형의 스타일 정보 수집
                for shape in slide.shapes:
                    shape_info = {
                        'type': shape.shape_type,
                        'name': shape.name
                    }
                    
                    # 텍스트 프레임이 있는 경우 텍스트 스타일 정보 추가
                    if shape.has_text_frame:
                        text_frame = shape.text_frame
                        paragraphs = []
                        for paragraph in text_frame.paragraphs:
                            paragraph_info = {
                                'text': paragraph.text,
                                'alignment': paragraph.alignment,
                                'font_name': paragraph.font.name if paragraph.font else None,
                                'font_size': paragraph.font.size if paragraph.font else None,
                                'font_bold': paragraph.font.bold if paragraph.font else None,
                                'font_italic': paragraph.font.italic if paragraph.font else None
                            }
                            paragraphs.append(paragraph_info)
                        shape_info['text_info'] = paragraphs
                    
                    slide_info['shapes'].append(shape_info)
                
                slides_info.append(slide_info)
            
            summary = {
                'slide_count': len(presentation.slides),
                'slides_info': slides_info,
                'core_properties': {
                    'author': presentation.core_properties.author,
                    'title': presentation.core_properties.title,
                    'subject': presentation.core_properties.subject,
                    'created': presentation.core_properties.created,
                    'modified': presentation.core_properties.modified
                }
            }
            
            return summary
        except Exception as e:
            logger.error(f"PowerPoint 데이터 요약 생성 중 오류 발생: {str(e)}")
            return None

    @staticmethod
    def extract_style_definitions(presentation: Presentation) -> Optional[Dict[str, Any]]:
        """
        PowerPoint 프레젠테이션의 스타일 정의를 추출합니다.
        
        Args:
            presentation (Presentation): PowerPoint 프레젠테이션 객체
            
        Returns:
            dict: 스타일 정의 정보
        """
        if presentation is None:
            return None
            
        try:
            # 마스터 슬라이드 스타일 정보
            master_styles = []
            for master in presentation.slide_masters:
                master_style = {
                    'name': master.name,
                    'layouts': []
                }
                
                # 레이아웃 스타일 정보
                for layout in master.slide_layouts:
                    layout_style = {
                        'name': layout.name,
                        'placeholders': []
                    }
                    
                    # 플레이스홀더 스타일 정보
                    for placeholder in layout.placeholders:
                        placeholder_style = {
                            'type': placeholder.placeholder_format.type if hasattr(placeholder, 'placeholder_format') else None,
                            'idx': placeholder.placeholder_format.idx if hasattr(placeholder, 'placeholder_format') else None
                        }
                        layout_style['placeholders'].append(placeholder_style)
                    
                    master_style['layouts'].append(layout_style)
                
                master_styles.append(master_style)
            
            # 테마 색상 정보
            theme_colors = {}
            if presentation.slide_masters[0].theme:
                theme = presentation.slide_masters[0].theme
                if theme.theme_elements.clrScheme:
                    for clr in theme.theme_elements.clrScheme:
                        theme_colors[clr.name] = clr.xml
            
            style_definitions = {
                'master_styles': master_styles,
                'theme_colors': theme_colors
            }
            
            return style_definitions
        except Exception as e:
            logger.error(f"PowerPoint 스타일 정의 추출 중 오류 발생: {str(e)}")
            return None

    @staticmethod
    def extract_content(presentation: Presentation) -> Optional[str]:
        """
        PowerPoint 프레젠테이션에서 텍스트 내용만 추출합니다.
        
        Args:
            presentation (Presentation): PowerPoint 프레젠테이션 객체
            
        Returns:
            str: 추출된 텍스트 내용
            None: 에러 발생 시
        """
        if presentation is None:
            return None
            
        try:
            content = []
            
            # 슬라이드 노트 포함 여부를 확인하고 처리
            has_notes = any(slide.has_notes_slide for slide in presentation.slides)
            
            for slide in presentation.slides:
                slide_texts = []
                
                # 슬라이드 제목 추출
                if slide.shapes.title:
                    slide_texts.append(slide.shapes.title.text)
                
                # 슬라이드의 모든 도형에서 텍스트 추출
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text.strip()
                        if text:
                            slide_texts.append(text)
                
                # 노트가 있는 경우 노트 내용도 추출
                if has_notes and slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slide_texts.append(f"Notes: {notes_text}")
                
                # 슬라이드의 모든 텍스트를 하나로 합침
                if slide_texts:
                    content.append("\n".join(slide_texts))
            
            # 모든 슬라이드의 내용을 하나의 문자열로 합침
            return "\n\n".join(content)
            
        except Exception as e:
            logger.error(f"PowerPoint 내용 추출 중 오류 발생: {str(e)}")
            return None 